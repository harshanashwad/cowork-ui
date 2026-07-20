"""
Copied into every session's working directory (see sessions.py) so the
agent can run it directly via bash — this is not just a backend-internal
module, it's the actual tool the agent uses.

Gives the agent structured slide geometry instead of it guessing overlap
by eye from the pptx file's raw internals — that's where the real risk
was (wrong-shape-targeting, guessed coordinates) before this existed.
"""

import json
import sys

from pptx import Presentation


def list_shape_bounds(pptx_path: str, slide_index: int) -> list[dict]:
    """
    Every existing shape's position and size on one slide, in inches
    (not raw EMU) so the numbers are directly usable when deciding where
    a new shape can go without overlapping anything already there.
    """
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]

    bounds = []
    for shape in slide.shapes:
        bounds.append({
            "name": shape.name,
            "type": str(shape.shape_type),
            "left_in": shape.left.inches if shape.left is not None else None,
            "top_in": shape.top.inches if shape.top is not None else None,
            "width_in": shape.width.inches if shape.width is not None else None,
            "height_in": shape.height.inches if shape.height is not None else None,
        })
    return bounds


if __name__ == "__main__":
    # Usage: python3 pptx_helpers.py deck_copy.pptx <slide_index (0-based)>
    path, slide_index = sys.argv[1], int(sys.argv[2])
    print(json.dumps(list_shape_bounds(path, slide_index), indent=2))
